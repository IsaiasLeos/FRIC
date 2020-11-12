import json
import pymongo
import random
import docx
import openpyxl
import pptx
from flask import Flask, jsonify, request, make_response
from docx import Document
from docx.shared import Inches, Pt
from datetime import date
from pptx import Presentation

from pptx.chart.data import CategoryChartData  
from pptx.enum.chart import XL_CHART_TYPE 


app = Flask(__name__)
#TO:DO 
# Add analyst
# Dont allow empty events
@app.route('/addAnalystToEvent',methods=['POST'])
def addAnalyst():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    myAnalystCollection = mydb["event_analyst"]
    req = request.get_json()

    analyst = {"event_id":req["id"], "analyst": req["analyst"], "is_lead": req["is_lead"]}
    myAnalystCollection.insert_one(analyst) 

    return "OK"



# Given Analyst return progress # tasks completed / # of tasks
def calculateProgress(analyst):
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    myTaskCollection = mydb["task"]
    tasks = []
    progress = 0
    for t in myTaskCollection.find({"Task_Analysts": analyst}):
        tasks.append(t)
        progress += int(t["Task_Progress"])
    if len(tasks) == 0:
        return 0
    return progress / len(tasks)


# Given event, return analysts from that event # 
@app.route('/analystsInEvent',methods=['POST'])
def analystList():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    analysts = []
    myAnalystCollection = mydb["event_analyst"]
    req = request.get_json()
    for a in myAnalystCollection.find({"event_id": req}):
        if a["is_lead"] == "0":
            analysts.append({"analyst": a["analyst"],"event":a["event_id"],"is_lead": a["is_lead"], "progress": calculateProgress(a["analyst"])})
    
    return jsonify(analysts)

# Given event, return lead analysts from that event # 
@app.route('/leadAnalystsInEvent',methods=['POST'])
def leadAnalystList():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    analysts = []
    myAnalystCollection = mydb["event_analyst"]
    req = request.get_json()
    for a in myAnalystCollection.find({"event_id": req}):
        if a["is_lead"] == "1":
            analysts.append({"analyst": a["analyst"],"event":a["event_id"],"is_lead": a["is_lead"], "progress": calculateProgress(a["analyst"])})
    
    return jsonify(analysts)


# Returns all analysts #
@app.route("/analysts")
def analysts():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    analysts = []
    myAnalystCollection = mydb["analyst"]

    for a in myAnalystCollection.find():
        analysts.append({"isLead": a["isLead"], "initials": a["initials"]})
    return jsonify(analysts)

#Return All events
@app.route('/eventsOverview')
def eventsOverview():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    myEventCollection = mydb["event"]
    mySystemCollection = mydb["system"]
    myFindingCollection = mydb["finding"]

    events_json = []

    findings_json = []
    # Get number of Findings
    for f in myFindingCollection.find():
        findings_json.append({"id": f["id"], "hostName": f["Host_Name"]})
    num_finds = len(findings_json)

    systems_json = []
    # Get number of systems
    for s in mySystemCollection.find():
        systems_json.append(
            {"sysInfo": s["System_Info"], "sysDesc": s["System_Description"]}
        )
    num_sys = len(systems_json)

    # Event Overview Information
    for e in myEventCollection.find():

        events_json.append(
            {
                "id": e["id"],
                "name": e["Event_name"],
                "desc": e["Description"],
                "type": e["Type"],
                "version": e["Version"],
                "assess_date": e["Assessment_date"],
                "org_name": e["Org_name"],
                "event_class": e["Event_class"],
                "declass_date": e["Declass_date"],
                "customer": e["Customer_name"],
                "num_sys": num_sys,
                "num_findings": num_finds,
                "prog": e["Progress"],
                "created_by": e["Created_By"],
            }
        )

    return jsonify(events_json)


# TO:DO Event ID Increment
@app.route("/addevent", methods=["POST"])
def addEvent():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    mycollection = mydb["event"]

    req = request.get_json()
    event = {
        "id": str(random.randint(1, 30)),
        "Event_name": req["name"],
        "Description": req["desc"],
        "Type": req["type"],
        "Version": req["vers"],
        "Assessment_date": req["assess_date"],
        "Org_name": req["org_name"],
        "Event_class": req["event_class"],
        "Declass_date": req["declass_date"],
        "Customer_name": req["customer_name"],
        "Created_By": req["created_by"],
        "Num_systems": 13,
        "Num_findings": 10,
        "Progress": "33%",
    }

    mycollection.insert_one(event)
    return


@app.route("/getprogress")
def getProgress():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    mySystemCollection = mydb["task"]
    task_progress = []
    for e in mySystemCollection.find():
        task_progress.append({"taskProgress": e["Task_Progress"]})
    return jsonify(task_progress)


@app.route("/getsystem")
def systems():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    mySystemCollection = mydb["system"]
    myFindingCollection = mydb["finding"]
    myTaskCollection = mydb["task"]
    system_json = []
    findings_json = []

    for f in myFindingCollection.find():
        findings_json.append({"id": f["id"], "hostName": f["Host_Name"]})
    num_finds = len(findings_json)

    task_json = []
    for f in myTaskCollection.find():
        task_json.append(
            {"taskTitle": f["Task_title"], "taskDescription": f["Task_Description"]}
        )
    num_tasks = len(task_json)

    for e in mySystemCollection.find():
        system_json.append(
            {
                "id": e["id"],
                "sysInfo": e["System_Info"],
                "sysDesc": e["System_Description"],
                "sysLoc": e["System_Location"],
                "sysRouter": e["System_Router"],
                "sysSwitch": e["System_Switch"],
                "sysRoom": e["System_Room"],
                "sysTestPlan": e["Test_Plan"],
                "Confidentiality": e["Confidentiality"],
                "Integrity": e["Integrity"],
                "Availability": e["Availability"],
                "num_task": num_tasks,
                "num_findings": num_finds,
                "prog": e["Progress"],
                "eventID": e["Event_ID"],
            }
        )
    return jsonify(system_json)


@app.route("/addsystem", methods=["POST"])
def addSystems():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    mycollection = mydb["system"]

    req = request.get_json()
    system = {
        "id": str(random.randint(1, 30)),
        "System_Info": req["sysInfo"],
        "System_Description": req["sysDesc"],
        "System_Location": req["sysLoc"],
        "System_Router": req["sysRouter"],
        "System_Switch": req["sysSwitch"],
        "System_Room": req["sysRoom"],
        "Test_Plan": req["sysTestPlan"],
        "Confidentiality": req["Confidentiality"],
        "Integrity": req["Integrity"],
        "Availability": req["Availability"],
        "Num_Task": 13,
        "Num_Findings": 10,
        "Progress": "0%",
        "Event_ID": req["eventID"],
    }
    mycollection.insert_one(system)
    return


@app.route("/editsystem", methods=["POST"])
def editSystem():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    mycollection = mydb["system"]
    req = request.get_json()
    query = {"id": req["id"]}
    system = {
        "$set": {
            "System_Info": req["sysInfo"],
            "System_Description": req["sysDesc"],
            "System_Location": req["sysLoc"],
            "System_Router": req["sysRouter"],
            "System_Switch": req["sysSwitch"],
            "System_Room": req["sysRoom"],
            "Test_Plan": req["sysTestPlan"],
            "Confidentiality": req["Confidentiality"],
            "Integrity": req["Integrity"],
            "Availability": req["Availability"],
            "Num_Task": 13,
            "Num_Findings": 10,
            "Progress": "0%",
            "Event_ID": req["eventID"],
        }
    }
    mycollection.update_one(query, system)
    return jsonify(system)


@app.route("/editevent", methods=["POST"])
def editEvent():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    mycollection = mydb["event"]

    req = request.get_json()
    query = {"id": req["id"]}
    event = {
        "$set": {
            "Event_name": req["name"],
            "Description": req["desc"],
            "Type": req["type"],
            "Version": req["vers"],
            "Assessment_date": req["assess_date"],
            "Org_name": req["org_name"],
            "Event_class": req["event_class"],
            "Declass_date": req["declass_date"],
            "Customer_name": req["customer_name"],
            "Created_By": req["created_by"],
            "Num_systems": 13,
            "Num_findings": 10,
            "Progress": "33%",
        }
    }
    mycollection.update_one(query, event)
    return jsonify(event)


# ---------------START OF FINDING API ---------------#

#Get the current analysts findings# 
@app.route("/analystFindings",methods=["POST"])  # path used in JS to call this
def analystFindings():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]  # Database name
    mycollection = mydb["finding"]  # Collection Name
    req = request.get_json()
    
    finding_json = []

    # Start of Finding
    for e in mycollection.find({"analyst": req}):
        
        finding_json.append(
            {
                "id": e["id"],
                "hostName": e["Host_Name"],
                "ip_port": e["IP_Port"],
                "description": e["Description"],
                "longDescription": e["Long_Description"],
                "findingStatus": e["Finding_Status"],
                "findingType": e["Finding_Type"],
                "findingClassification": e["Finding_Classification"],
                "findingSystem": e["Finding_System"],
                "findingTask": e["Finding_Task"],
                "findingSubtask": e["Finding_Subtask"],
                "relatedFindings": e["Related_Findings"],
                "findingConfidentiality": e["Finding_Confidentiality"],
                "findingIntegrity": e["Finding_Integrity"],
                "findingAvailability": e["Finding_Availability"],
                "findingAnalyst": e["Finding_Analyst"],
                "findingCollaborators": e["Finding_Collaborators"],
                "findingPosture": e["Finding_Posture"],
                "mitigationDesc": e["Mitigation_Desc"],
                "mitigationLongDesc": e["Mitigation_Long_Desc"],
                "threatRelevence": e["Threat_Relevence"],
                "countermeasure": e["Countermeasure"],
                "impactDesc": e["Impact_Desc"],
                "impactLevel": e["Impact_Level"],
                "severityCategoryScore": e["Severity_Score"],
                "vulnerabilityScore": e["Vulnerability_Score"],
                "quantitativeScore": e["Quantitative_Score"],
                "findingRisk": e["Finding_Risk"],
                "findingLikelihood": e["Finding_Likelihood"],
                "findingCFIS": e["Finding_CFIS"],
                "findingIFIS": e["Finding_IFIS"],
                "findingAFIS": e["Finding_AFIS"],
                "impactScore": e["Impact_Score"],
                "findingFiles": e["Finding_Files"],
                "severityCategoryCode": e["Severity_Category_Code"],
                "systemID": e["System_ID"],
                "taskID": e["Task_ID"],
                "subtaskID": e["Subtask_ID"],
            }
        )
        
    return jsonify(finding_json)  # return what was found in the collection


@app.route("/findings")  # path used in JS to call this
def findings():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]  # Database name
    mycollection = mydb["finding"]  # Collection Name

    finding_json = []

    # Start of Finding
    for e in mycollection.find():
        finding_json.append(
            {
                "id": e["id"],
                "hostName": e["Host_Name"],
                "ip_port": e["IP_Port"],
                "description": e["Description"],
                "longDescription": e["Long_Description"],
                "findingStatus": e["Finding_Status"],
                "findingType": e["Finding_Type"],
                "findingClassification": e["Finding_Classification"],
                "findingSystem": e["Finding_System"],
                "findingTask": e["Finding_Task"],
                "findingSubtask": e["Finding_Subtask"],
                "relatedFindings": e["Related_Findings"],
                "findingConfidentiality": e["Finding_Confidentiality"],
                "findingIntegrity": e["Finding_Integrity"],
                "findingAvailability": e["Finding_Availability"],
                "findingAnalyst": e["Finding_Analyst"],
                "findingCollaborators": e["Finding_Collaborators"],
                "findingPosture": e["Finding_Posture"],
                "mitigationDesc": e["Mitigation_Desc"],
                "mitigationLongDesc": e["Mitigation_Long_Desc"],
                "threatRelevence": e["Threat_Relevence"],
                "countermeasure": e["Countermeasure"],
                "impactDesc": e["Impact_Desc"],
                "impactLevel": e["Impact_Level"],
                "severityCategoryScore": e["Severity_Score"],
                "vulnerabilityScore": e["Vulnerability_Score"],
                "quantitativeScore": e["Quantitative_Score"],
                "findingRisk": e["Finding_Risk"],
                "findingLikelihood": e["Finding_Likelihood"],
                "findingCFIS": e["Finding_CFIS"],
                "findingIFIS": e["Finding_IFIS"],
                "findingAFIS": e["Finding_AFIS"],
                "impactScore": e["Impact_Score"],
                "findingFiles": e["Finding_Files"],
                "severityCategoryCode": e["Severity_Category_Code"],
                "systemID": e["System_ID"],
                "taskID": e["Task_ID"],
                "subtaskID": e["Subtask_ID"]
            }
        )
    return jsonify(finding_json)  # return what was found in the collection


# ---------- HELPER FUNCTIONS TO DERIVE ATTRIBUTES OF FINDING ----------#

# Map the Finding Severity Code to the Score
def calculateSeverityScore(code):
    severityCategoryScore = 0

    if code == "I":
        severityCategoryScore = 10
    elif code == "II":
        severityCategoryScore = 7
    else:
        severityCategoryScore = 4

    return severityCategoryScore


# Function to calculate impact score based on the finding CFIS IFIS and AFIS
def calculateImpactScore(CFIS, IFIS, AFIS):
    findingSystemLevel = ""
    impact_score = 0
    findingSystemLevel += CFIS  # Create combination string of the values
    findingSystemLevel += IFIS
    findingSystemLevel += AFIS

    systemlevelQuantitative = {  # Map combo to value given in SRS
        "HHH": 10,
        "HHX": 9,
        "HXX": 8,
        "MMM": 7,
        "MMX": 6,
        "MXX": 5,
        "LLL": 4,
        "LLX": 3,
        "LXX": 2,
        "XXX": 0,
    }

    impact_score = systemlevelQuantitative.get(findingSystemLevel)
    return impact_score


# Function to calculate Vulnerability Severity based on countermeasure, impactscore, and severity category score
def calculateVulnerabilitySeverity(countermeasure, impactScore, severityCategoryScore):
    vulnerabilitySeverityScore = 0
    countermeasureScore = 0

    countermeasureScore = int(countermeasure)
    vulnerabilitySeverityScore = (
        countermeasureScore * impactScore * severityCategoryScore
    ) / 10  # Algorithm used in SRS to derive Severity Score

    return vulnerabilitySeverityScore


# Function to calculate QVS based on the vulnerability severity score
def calcualteQuantitativeVulnerabilitySeverity(vulnerabilitySeverityScore):
    quantitativeVulnerabilitySeverityScore = ""

    if (
        vulnerabilitySeverityScore >= 95 and vulnerabilitySeverityScore <= 100
    ):  # Assignment based on value
        quantitativeVulnerabilitySeverityScore = "VH"
    elif vulnerabilitySeverityScore >= 80 and vulnerabilitySeverityScore < 95:
        quantitativeVulnerabilitySeverityScore = "H"
    elif vulnerabilitySeverityScore >= 20 and vulnerabilitySeverityScore < 80:
        quantitativeVulnerabilitySeverityScore = "M"
    elif vulnerabilitySeverityScore >= 5 and vulnerabilitySeverityScore < 20:
        quantitativeVulnerabilitySeverityScore = "L"
    else:
        quantitativeVulnerabilitySeverityScore = "VL"

    return quantitativeVulnerabilitySeverityScore


# Function to assign index (for future mapping) based on the value of impact (FOR FINDING)
def routeImpact(impact):
    impactIndices = {  # Mapping of possible index
        "VL": 0,
        "L": 1,
        "M": 2,
        "H": 3,
        "VH": 4,
    }
    impact_index = impactIndices.get(impact)
    return impact_index


# Function to assign index (for future mapping) based on the value of likelihood
def routeLikelihood(likelihood):
    likelihoodIndices = {  # Mapping of possible index
        "VH": 0,
        "H": 1,
        "M": 2,
        "L": 3,
        "VL": 4,
    }
    likelihood_index = likelihoodIndices.get(likelihood)
    return likelihood_index


# Function to assign index (for future mapping) based on the value of severity
def routeVulnerabilitySeverity(severity):
    vulnerabilityIndices = {  # Mapping of possible index
        "VL": 0,
        "L": 1,
        "M": 2,
        "H": 3,
        "VH": 4,
        "INFO": 5,
    }
    vulnerability_index = vulnerabilityIndices.get(severity)
    return vulnerability_index


# Function to assign index (for future mapping) based on the value of threat
def routeRelevenceOfThreat(threat):
    threatIndices = {  # Map index based on value
        "Confirmed": 0,
        "Expected": 1,
        "Anticipated": 2,
        "Predicted": 3,
        "Possible": 4,
    }
    threat_index = threatIndices.get(threat)
    return threat_index


# Function to calculate the Finding Likelihood
def calculateLikelihood(relevenceOfThreat, vulnerabilitySeverity):
    likelihoodMap = [  # Pre defined map per SRS
        ["VL", "L", "M", "H", "VH"],
        ["VL", "L", "M", "H", "VH"],
        ["VL", "L", "M", "M", "H"],
        ["VL", "L", "L", "L", "M"],
        ["VL", "VL", "L", "L", "L"],
    ]

    threat = routeRelevenceOfThreat(relevenceOfThreat)  # Get index
    vulnerability = routeVulnerabilitySeverity(vulnerabilitySeverity)  # Get index

    likelihood = likelihoodMap[threat][vulnerability]  # Select value based on indices
    return likelihood


# Function to calculate the Finding Risk
def calculateRisk(likelihood, impactLevel):
    riskMap = [  # Pre defined map per SRS
        ["VL", "L", "M", "H", "VH"],
        ["VL", "L", "M", "H", "VH"],
        ["VL", "L", "M", "M", "H"],
        ["VL", "L", "L", "L", "M"],
        ["VL", "VL", "L", "L", "L"],
        ["INFO", "INFO", "INFO", "INFO", "INFO"],
    ]

    impact = routeImpact(impactLevel)  # Get index for impact ('Y' value)
    Likelihood = routeLikelihood(likelihood)  # Get index for Likelihood ('X' value)

    risk = riskMap[Likelihood][impact]  # Select value based on indices
    return risk


@app.route("/addfinding", methods=["POST"])
def addFindings():
    myclient = pymongo.MongoClient(
        "mongodb://localhost:27017/"
    )  # Connect to the DB Client
    mydb = myclient["FRIC"]
    mycollection = mydb["finding"]

    req = request.get_json()

    # severityCategoryScore = 0 #Derived from Severity Category Code

    finding = {
        "id": str(random.randint(1, 30)),
        "Host_Name": req["hostName"],
        "IP_Port": req["ip_port"],
        "Description": req["description"],
        "Long_Description": req["longDescription"],
        "Finding_Status": req["findingStatus"],
        "Finding_Type": req["findingType"],
        "Finding_Classification": req["findingClassification"],
        "Finding_System": req["findingSystem"],
        "Finding_Task": req["findingTask"],
        "Finding_Subtask": req["findingSubtask"],
        "Related_Findings": req["relatedFindings"],
        "Finding_Confidentiality": req["findingConfidentiality"],
        "Finding_Integrity": req["findingIntegrity"],
        "Finding_Availability": req["findingAvailability"],
        "Finding_Analyst": req["findingAnalyst"],
        "Finding_Collaborators": req["findingCollaborators"],
        "Finding_Posture": req["findingPosture"],
        "Mitigation_Desc": req["mitigationDesc"],
        "Mitigation_Long_Desc": req["mitigationLongDesc"],
        "Threat_Relevence": req["threatRelevence"],
        "Countermeasure": req["countermeasure"],
        "Impact_Desc": req["impactDesc"],
        "Impact_Level": req["impactLevel"],
        "Severity_Score": req["severityCategoryScore"],
        "Vulnerability_Score": req["vulnerabilityScore"],
        "Quantitative_Score": req["quantitativeScore"],
        "Finding_Risk": req["findingRisk"],
        "Finding_Likelihood": req["findingLikelihood"],
        "Finding_CFIS": req["findingCFIS"],
        "Finding_IFIS": req["findingIFIS"],
        "Finding_AFIS": req["findingAFIS"],
        "Impact_Score": req["impactScore"],
        "Finding_Files": req["findingFiles"],
        "Severity_Category_Code": req["severityCategoryCode"],
        "System_ID": req["systemID"],
        "Task_ID": req["taskID"],
        "Subtask_ID": req["subtaskID"],
        "analyst": req["analyst"]
    }

    # ----START OF DERIVED ATTRIBUTES----#
    # Calculate Severity Category Score
    severityCategoryScore = 0
    severityCategoryCode = finding.get("Severity_Category_Code")
    severityCategoryScore = calculateSeverityScore(severityCategoryCode)
    finding.update({"Severity_Score": severityCategoryScore})

    # Calculate Impact Score
    findingImpactScore = 0
    findingCFIS = finding.get("Finding_CFIS")
    findingIFIS = finding.get("Finding_IFIS")
    findingAFIS = finding.get("Finding_AFIS")
    findingImpactScore = calculateImpactScore(findingCFIS, findingIFIS, findingAFIS)
    finding.update({"Impact_Score": findingImpactScore})

    # Calculate Vulerability Severity=
    vulnerabilitySeverityScore = 0
    counterMeasure = finding.get("Countermeasure")
    vulnerabilitySeverityScore = calculateVulnerabilitySeverity(
        counterMeasure, findingImpactScore, severityCategoryScore
    )
    finding.update({"Vulnerability_Score": vulnerabilitySeverityScore})

    # Calculate Quantitative Vulnerability Severity
    QVS = ""
    QVS = calcualteQuantitativeVulnerabilitySeverity(vulnerabilitySeverityScore)
    finding.update({"Quantitative_Score": QVS})

    # Calculate Likelihood
    threat_relevence = ""
    threat_relevence = finding.get("Threat_Relevence")
    likelihood = ""

    if findingImpactScore == 0:
        likelihood = "INFO"
    else:
        likelihood = calculateLikelihood(threat_relevence, QVS)

    finding.update({"Finding_Likelihood": likelihood})

    # Calculate Risk
    impact_level = ""
    impact_level = finding.get("Impact_Level")
    risk = ""

    if findingImpactScore == 0:
        risk = "INFO"
    else:
        risk = calculateRisk(likelihood, impact_level)

    finding.update({"Finding_Risk": risk})
    # ----END OF DERIVED ATTRIBUTES----#

    mycollection.insert_one(finding)  # Send information to collection
    return "OK"


@app.route("/editfinding", methods=["POST"])
def editFinding():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    mycollection = mydb["finding"]
    finding = []

    req = request.get_json()

    query = {"id": req["id"]}

    finding = {
        "$set": {
            "Host_Name": req["hostName"],
            "IP_Port": req["ip_port"],
            "Description": req["description"],
            "Long_Description": req["longDescription"],
            "Finding_Status": req["findingStatus"],
            "Finding_Type": req["findingType"],
            "Finding_Classification": req["findingClassification"],
            "Finding_System": req["findingSystem"],
            "Finding_Task": req["findingTask"],
            "Finding_Subtask": req["findingSubtask"],
            "Related_Findings": req["relatedFindings"],
            "Finding_Confidentiality": req["findingConfidentiality"],
            "Finding_Integrity": req["findingIntegrity"],
            "Finding_Availability": req["findingAvailability"],
            "Finding_Analyst": req["findingAnalyst"],
            "Finding_Collaborators": req["findingCollaborators"],
            "Finding_Posture": req["findingPosture"],
            "Mitigation_Desc": req["mitigationDesc"],
            "Mitigation_Long_Desc": req["mitigationLongDesc"],
            "Threat_Relevence": req["threatRelevence"],
            "Countermeasure": req["countermeasure"],
            "Impact_Desc": req["impactDesc"],
            "Impact_Level": req["impactLevel"],
            "Severity_Score": req["severityCategoryScore"],
            "Vulnerability_Score": req["vulnerabilityScore"],
            "Quantitative_Score": req["quantitativeScore"],
            "Finding_Risk": req["findingRisk"],
            "Finding_Likelihood": req["findingLikelihood"],
            "Finding_CFIS": req["findingCFIS"],
            "Finding_IFIS": req["findingIFIS"],
            "Finding_AFIS": req["findingAFIS"],
            "Impact_Score": req["impactScore"],
            "Finding_Files": req["findingFiles"],
            "Severity_Category_Code": req["severityCategoryCode"],
            "System_ID": req["systemID"],
            "Task_ID": req["taskID"],
            "Subtask_ID": req["subtaskID"],
        }
    }

    mycollection.update_one(query, finding)

    return jsonify(finding)


# --------------- END OF FINDING API ---------------#


# ---------------START OF SUBTASK API ---------------#
@app.route("/subtasks")
def subtasks():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    mycollection = mydb["subtask"]
    myFindingCollection = mydb["finding"]
    subtask_json = []

    findings_json = []
    # Get number of Findings
    for f in myFindingCollection.find():
        findings_json.append({"id": f["id"], "hostName": f["Host_Name"]})
    num_finds = len(findings_json)

    for e in mycollection.find():
        subtask_json.append(
            {
                "id": e["id"],
                "subtaskTitle": e["Subtask_Title"],
                "subtaskDescription": e["Subtask_Description"],
                "subtaskProgress": e["Subtask_Progress"],
                "subtaskDueDate": e["Subtask_Due_Date"],
                "analysts": e["Analysts"],
                "collaborators": e["Collaborators"],
                "relatedTask": e["Related_Task"],
                "subtasks": e["Subtasks"],
                "attachments": e["Attachments"],
                "numFindings": num_finds,
                "analyst": e["Analyst"],
                "task": e["Task"],
                "taskID": e["Task_ID"],
            }
        )
    return jsonify(subtask_json)


@app.route("/addsubtask", methods=["POST"])
def addSubtasks():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    mycollection = mydb["subtask"]
    req = request.get_json()
    subtask = {
        "id": str(random.randint(1, 30)),
        "Subtask_Title": req["subtaskTitle"],
        "Subtask_Description": req["subtaskDescription"],
        "Subtask_Progress": req["subtaskProgress"],
        "Subtask_Due_Date": req["subtaskDueDate"],
        "Analysts": req["analysts"],
        "Collaborators": req["collaborators"],
        "Related_Task": req["relatedTask"],
        "Subtasks": req["subtasks"],
        "Attachments": req["attachments"],
        "Num_Findings": 0,
        "Analyst": "Analyst 0",
        "Task": "Task 0",
        "Task_ID": req["taskID"],
    }
    mycollection.insert_one(subtask)


@app.route("/editsubtask", methods=["POST"])
def editSubtask():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    mycollection = mydb["subtask"]
    req = request.get_json()
    query = {"id": req["id"]}
    subtask = {
        "$set": {
            "Subtask_Title": req["subtaskTitle"],
            "Subtask_Description": req["subtaskDescription"],
            "Subtask_Progress": req["subtaskProgress"],
            "Subtask_Due_Date": req["subtaskDueDate"],
            "Analysts": req["analysts"],
            "Collaborators": req["collaborators"],
            "Related_Task": req["relatedTask"],
            "Subtasks": req["subtasks"],
            "Attachments": req["attachments"],
            "Num_Findings": 0,
            "Analyst": "Analyst 0",
            "Task": "Task 0",
            "Task_ID": req["taskID"],
        }
    }
    mycollection.update_one(query, subtask)
    return jsonify(subtask)


# --------------- END OF SUBTASK API ---------------#


# ---------------START OF TASK API ---------------#
# Function used to get overview of tasks
@app.route("/tasks")
def tasks():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient['FRIC']
    mycollection = mydb['task']
    myFindingCollection = mydb['finding']
    mySubtaskCollection = mydb['subtask']
    
    findings_json = []
    # Get number of Findings
    for f in myFindingCollection.find():
        findings_json.append({"id": f["id"], "hostName": f["Host_Name"]})
    num_finds = len(findings_json)

    subtask_json = []
    # Get number of systems
    for s in mySubtaskCollection.find():
        subtask_json.append(
            {
                "subtaskTitle": s["Subtask_Title"],
                "subtaskDescription": s["Subtask_Description"],
            }
        )
    num_subtask = len(subtask_json)

    task_json = []
    # Start of task 
    for e in mycollection.find():

        task_json.append(
            {
                "id": e["id"],
                "taskTitle": e["Task_title"],
                "taskDescription": e["Task_Description"],
                "system": e["System"],
                "taskPriority": e["Task_Priority"],
                "taskProgress": e["Task_Progress"],
                "taskDueDate": e["Task_Due_Date"],
                "taskAnalysts": e["Task_Analysts"],
                "taskCollaborators": e["Task_Collaborators"],
                "relatedTasks": e["Related_Tasks"],
                "attachments": e["Attachments"],
                "num_subtask": num_subtask,
                "num_finding": num_finds,
                "subtaskID": e["SubTask_ID"],
            }
        )
    return jsonify(task_json)


# Function used to add task
@app.route("/addtask", methods=["POST"])
def addTasks():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    mycollection = mydb["task"]
    req = request.get_json()
    task = {
        "id": str(random.randint(1, 30)),
        "Task_title": req["taskTitle"],
        "Task_Description": req["taskDescription"],
        "System": req["system"],
        "Task_Priority": req["taskPriority"],
        "Task_Progress": req["taskProgress"],
        "Task_Due_Date": req["taskDueDate"],
        "Task_Analysts": req["taskAnalysts"],
        "Task_Collaborators": req["taskCollaborators"],
        "Related_Tasks": req["relatedTasks"],
        "Attachments": req["attachments"],
        "Num_subtask": 0,
        "Num_finding": 13,
        "Progress": "0%",
        "SubTask_ID": req["subtaskID"],
    }
    mycollection.insert_one(task)  # send info to collection
    return "OK"


# Function used to edit task
@app.route("/edittask", methods=["POST"])
def editTask():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    mycollection = mydb["task"]

    req = request.get_json()
    query = {"id": req["id"]}

    task = {
        "$set": {
            "Task_title": req["taskTitle"],
            "Task_Description": req["taskDescription"],
            "System": req["system"],
            "Task_Priority": req["taskPriority"],
            "Task_Progress": req["taskProgress"],
            "Task_Due_Date": req["taskDueDate"],
            "Task_Analysts": req["taskAnalysts"],
            "Task_Collaborators": req["taskCollaborators"],
            "Related_Tasks": req["relatedTasks"],
            "Attachments": req["attachments"],
            "Num_subtask": 0,
            "Num_finding": 13,
            "Progress": "0%",
            "SubTask_ID": req["subtaskID"],
        }
    }
    mycollection.update_one(query, task)
    return jsonify(task)


# --------------- END OF TASK API ---------------#


@app.route("/addlog", methods=["POST"])
def addLog():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    mycollection = mydb["logs"]

    req = request.get_json()
    log = {
        "Date_Time": req["date"],
        "Action_Performed": req["action"],
        "Analyst": req["analyst"],
    }
    mycollection.insert_one(log)
    return "OK"

@app.route('/createRiskMatrix', methods=['POST'])
def create_Risk_Matrix():
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient['FRIC']
    myFindingCollection = mydb['finding']
    finding_json = []
    for f in myFindingCollection.find():    #Getting the findings in the db
        finding_json.append(
            (f["IP_Port"], f["Description"], f["Finding_Status"], f["Finding_Type"], f["Finding_Posture"], f["Finding_Confidentiality"], f["Finding_Integrity"], f["Finding_Availability"], f["Impact_Score"], f["Severity_Category_Code"], f["Severity_Score"], f["Countermeasure"], f["Vulnerability_Score"], f["Quantitative_Score"], f["Threat_Relevence"], f["Finding_Likelihood"], f["Impact_Level"], f["Finding_Risk"] ))
    wb = openpyxl.Workbook()    #Opening the workbook
    ws = wb.active  # Worksheet object
    ws.title = "risk_matrix"    # Changing the title of the worksheet
    ws.append(("IP:PORT","DESCRIPTION", "STATUS", "TYPE",  "POSTURE", "C","I","A", "IMP. SCORE", "CAT", "CAT SCORE", "CM", "VS(n)", "VS(q)" , "RELEVANCE OF THREAT", "LIKELIHOOD", "IMPACT", "RISK"))   #First row in the worksheet
    for finding in finding_json:        #Appending all of the findings to the worksheet
        ws.append(finding)
    wb.save("riskMatrix.xlsx")  #Saving the file

@app.route("/generateERB" , methods=["POST"])
def generateERB():
    
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]
    myEventCollection = mydb["event"]
    mySystemCollection = mydb["system"]
    myFindingCollection = mydb["finding"]
    
    events_json = []
    system_json = []
    finding_json = []

    # Start of Finding
    for e in myFindingCollection.find():
        finding_json.append(
            {
                "id": e["id"],
                "hostName": e["Host_Name"],
                "ip_port": e["IP_Port"],
                "description": e["Description"],
                "longDescription": e["Long_Description"],
                "findingStatus": e["Finding_Status"],
                "findingType": e["Finding_Type"],
                "findingClassification": e["Finding_Classification"],
                "findingSystem": e["Finding_System"],
                "findingTask": e["Finding_Task"],
                "findingSubtask": e["Finding_Subtask"],
                "relatedFindings": e["Related_Findings"],
                "findingConfidentiality": e["Finding_Confidentiality"],
                "findingIntegrity": e["Finding_Integrity"],
                "findingAvailability": e["Finding_Availability"],
                "findingAnalyst": e["Finding_Analyst"],
                "findingCollaborators": e["Finding_Collaborators"],
                "findingPosture": e["Finding_Posture"],
                "mitigationDesc": e["Mitigation_Desc"],
                "mitigationLongDesc": e["Mitigation_Long_Desc"],
                "threatRelevence": e["Threat_Relevence"],
                "countermeasure": e["Countermeasure"],
                "impactDesc": e["Impact_Desc"],
                "impactLevel": e["Impact_Level"],
                "severityCategoryScore": e["Severity_Score"],
                "vulnerabilityScore": e["Vulnerability_Score"],
                "quantitativeScore": e["Quantitative_Score"],
                "findingRisk": e["Finding_Risk"],
                "findingLikelihood": e["Finding_Likelihood"],
                "findingCFIS": e["Finding_CFIS"],
                "findingIFIS": e["Finding_IFIS"],
                "findingAFIS": e["Finding_AFIS"],
                "impactScore": e["Impact_Score"],
                "findingFiles": e["Finding_Files"],
                "severityCategoryCode": e["Severity_Category_Code"],
                "systemID": e["System_ID"],
                "taskID": e["Task_ID"],
                "subtaskID": e["Subtask_ID"],
            }
        )
   
    for e in myEventCollection.find():

        events_json.append(
            {
            
                "name": e["Event_name"],
                
            }
        )

    eventName = '' #Hold Event Name
    for x in range(len(events_json)): #Get the name of the event 
        event = events_json[x]
        eventName = event["name"]

    for e in mySystemCollection.find():
        system_json.append(
            {
                "sysInfo": e["System_Info"],
                
            }
        )

    ppt = Presentation()
    
    img_path = "../src/assets/logo.png"
    img_path2= "../src/assets/armyLogo.png"


    blank_slide_layout = ppt.slide_layouts[6]  
    
    # Attaching slide to ppt 
    slide = ppt.slides.add_slide(blank_slide_layout)  
    
    # For margins 
    left = Inches(.5)
    top = Inches(0)    
    height = Inches(1)  
    
    #Add image
    pic = slide.shapes.add_picture(img_path, left, top, height = height)

    left = Inches(8)
    top = Inches(0)    
    height = Inches(1.4)  
    
    pic = slide.shapes.add_picture(img_path2, left, top, height = height)

    # Client Information 
    left= Inches(0)
    top = Inches(1)
    height = Inches(1) 
    width = Inches(6)
    txtBox = slide.shapes.add_textbox(left,top ,width,height)
    tf = txtBox.text_frame
    tf.text = ""
    p = tf.add_paragraph()
    p.text = "U.S. ARMY COMBAT CAPABILITIES DEVELOPMENT COMMAND - DATA & ANALYSIS CENTER"
    
    p.font.size = Pt(19)
    p.font.bold = True


     #Client Information
    clientTxtBox = slide.shapes.add_textbox(Inches(3), Inches(2), width, height)
    tf4 = clientTxtBox.text_frame
    tf4.text = ""
    p4 = tf4.add_paragraph()
    p4.text = "Cyber Experimentation & Analysis Division"
    p4.font.size = Pt(17)
    p4.font.bold = True 

    #Event Information
    eventTxtBox = slide.shapes.add_textbox(Inches(4.5), Inches(3), width, height)
    tf2 = eventTxtBox.text_frame
    tf2.text = ""
    p2 = tf2.add_paragraph()
    p2.text = eventName #Event that took place
    p2.font.size = Pt(30)
    p2.font.bold = True    

    #Presenter Information
    presenterTxtBox = slide.shapes.add_textbox(Inches(0), Inches(6), width, height)
    tf3 = presenterTxtBox.text_frame
    tf3.text = "Name of Presenter: "
    p3 = tf3.add_paragraph()
    p3.text = "Rank/Title of Presenter: "
    p3.font.size = Pt(15)
    p3.font.bold = True 
    p5 = tf3.add_paragraph()

    today = "DD/MM/YYYY"
    p5.text = today
    p5.font.size = Pt(15)
    p5.font.bold = True 

    #Slide 1 Done
    
    #Slide 2
    blank_slide_layout2 = ppt.slide_layouts[6]  #Layout
    slide2 = ppt.slides.add_slide(blank_slide_layout2) #add to ppt

    # For margins 
    left = Inches(.5)
    top = Inches(0)    
    height = Inches(1)  
    
    #Add image
    pic = slide2.shapes.add_picture(img_path, left, top, height = height)

    left = Inches(8)
    top = Inches(0)    
    height = Inches(1.4)  
    
    pic = slide2.shapes.add_picture(img_path2, left, top, height = height)


    left= Inches(0)
    top = Inches(1)
    height = Inches(1) 
    width = Inches(6)
    txtBox = slide2.shapes.add_textbox(left,top ,width,height)
    tf = txtBox.text_frame
    tf.text = ""
    p = tf.add_paragraph()
    p.text = "Scope:"
    
    p.font.size = Pt(40)
    p.font.bold = True

    left= Inches(0)
    top = Inches(2)
    height = Inches(1) 
    width = Inches(6)
    sysTextBox = slide2.shapes.add_textbox(left,top ,width,height)
    tfA = sysTextBox.text_frame #Text Frame A
    tfA.level = 0
    tfA.text = "Systems assessed during the CVPA are as follows:"
    
    for x in range(len(system_json)): #List all systems in the given event
        pA = tfA.add_paragraph() #Paragraph A
        
        system = system_json[x]
        pA.level = 1
        pA.text = system["sysInfo"]
        pA.font.size = Pt(30)
        
    #----- END SLIDE 2 -----#
    
    #----- Start Slide 3 -----#

    slideTable = ppt.slides.add_slide(ppt.slide_layouts[5])
    x, y, cx, cy = Inches(0), Inches(2), Inches(10), Inches(4)
    shape = slideTable.shapes.add_table(len(finding_json)+1,5,x,y,cx,cy)
    table = shape.table

    cellID = table.cell(0,0)
    cellID.text = "ID"

    cellSystem = table.cell(0,1)
    cellSystem.text = "System"

    cellFinding = table.cell(0,2)
    cellFinding.text = "Finding"

    cellImpact = table.cell(0,3)
    cellImpact.text = "Impact"

    cellRisk = table.cell(0,4)
    cellRisk.text = "Risk"

    for x in range(1,len(finding_json) + 1):
        finding = finding_json[x-1] #Cant start at index zero because that is where labels are, however we still need the first finding to be put on the table

        finID = table.cell(x,0)
        finID.text = finding["id"]

        finSys = table.cell(x,1)
        finSys.text = finding["systemID"]

        finHostName = table.cell(x,2)
        finHostName.text = finding["hostName"]

        finImpact = table.cell(x,3)
        finImpact.text = finding["impactLevel"]

        finRisk = table.cell(x,4)
        finRisk.text = finding["findingRisk"]

    #Table Done
        
    left = Inches(.5)
    top = Inches(0)    
    height = Inches(1)  
    
    #Add image
    pic = slideTable.shapes.add_picture(img_path, left, top, height = height)

    left = Inches(8)
    top = Inches(0)    
    height = Inches(1.4)  
    
    pic = slideTable.shapes.add_picture(img_path2, left, top, height = height)
    left= Inches(0)
    top = Inches(1)
    height = Inches(1) 
    width = Inches(6)
    findingTxtBox = slideTable.shapes.add_textbox(left,top ,width,height)
    findingTf = findingTxtBox.text_frame
    findingTf.text = ""
    findingParagraph = findingTf.add_paragraph()
    findingParagraph.text = "Findings:"
    findingParagraph.font.size= Pt(30)

    #----- END SlIDE 3-----# 
    
    #----- START FINDINGS TABLES -----#
    for x in range(len(finding_json)):
        finding = finding_json[x]
        slideFinding= ppt.slides.add_slide(ppt.slide_layouts[5])

        x, y, cx, cy = Inches(0), Inches(2), Inches(10), Inches(4)

        shape = slideFinding.shapes.add_table(10,7,x,y,cx,cy)
        table = shape.table

        cellID = table.cell(0,0)
        cellID.text = "ID"
        cellA = table.cell(0,1)
        cellA.text = str(finding["id"])

        cellB = table.cell(0,2)
        cellB.text = "Impact Score"
        cellC = table.cell(0,3)
        cellC.text = str(finding['impactScore'])

        cellD = table.cell(0,4)
        cellD.text = "Status"
        cellE = table.cell(0,5)
        cellE.text = str(finding['findingStatus'])

        cellF = table.cell(0,6)
        cellF.text = 'Posture'
        cellG = table.cell(1,6)
        cellG.text = str(finding['findingPosture'])
        #Row 2

        cellH = table.cell(1,0)
        cellH.text = "Host Names"
        cellI= table.cell(1,1)
        cellI.text = str(finding["hostName"])

        cellJ = table.cell(2,0)
        cellJ.text = "IP:PORT"
        cellK = table.cell(2,1)
        cellK.text = str(finding['ip_port'])

        cellL = table.cell(1,2)
        cellL.text = "CAT"

        cellM = table.cell(1,3)
        cellM.text = str(finding['severityCategoryCode'])

        cellN= table.cell(1,4)
        cellN.text = "Likelihood"

        cellO = table.cell(1,5)
        cellO.text = str(finding['findingLikelihood'])

        cellP = table.cell(2,2)
        cellP.text = "CAT Score"

        cellq = table.cell(2,3)
        cellq.text = str(finding['severityCategoryScore'])

        cellR = table.cell(2,4)
        cellR.text = "Impact"

        cellS = table.cell(2,5)
        cellS.text = str(finding['impactLevel'])

        cellT = table.cell(3,0)
        cellT.text = 'Vs Score'

        cellU = table.cell(3,1)
        cellU.text = str(finding['vulnerabilityScore'])

        cellV = table.cell(3,2)
        cellV.text = "Risk"

        cellW = table.cell(3,3)
        cellW.text = str(finding['findingRisk'])

        cellX = table.cell(3,4)
        cellX.text = "Vs"

        cellY = table.cell(3,5)
        cellY.text = str(finding['vulnerabilityScore'])

        cellZ = table.cell(4,0)
        cellZ.text = 'CM'

        cellA1 = table.cell(4,1)
        cellA1.text = str(finding['countermeasure'])

        cellA2 = table.cell(4,2)
        cellA2.text = 'C:'

        cellA3 = table.cell(4,3)
        cellA3.text = str(finding['findingCFIS'])

        cellA4 = table.cell(4,4)
        cellA4.text = 'I:'

        cellA5 = table.cell(4,5)
        cellA5.text = str(finding['findingIFIS'])

        cellA6 = table.cell(5,0)
        cellA6.text = 'A:'

        cellA7 = table.cell(5,1)
        cellA7.text = str(finding['findingAFIS'])

        cellA8 = table.cell(5,2)
        cellA8.text = 'Impact Rationale'

        cellA9 = table.cell(5,3)
        cellA9.text = str(finding['impactDesc'])

        cellA10= table.cell(5,4)
        cellA10.text = 'Finding Type'

        cellA11 = table.cell(5,5)
        cellA11.text = str(finding['findingType'])

        cellA12 = table.cell(6,0)
        cellA12.text = 'Description'

        cellA13 = table.cell(6,1)
        cellA13.text = str(finding['longDescription'])

        cellA14 = table.cell(7,0)
        cellA14.text = 'Mitigation'

        cellA15 = table.cell(7,1)
        cellA15.text = str(finding['mitigationLongDesc'])

        #Add image
        left = Inches(.5)
        top = Inches(0)    
        height = Inches(1)  
        pic = slideFinding.shapes.add_picture(img_path, left, top, height = height)

        left = Inches(8)
        top = Inches(0)    
        height = Inches(1.4)  
        pic = slideFinding.shapes.add_picture(img_path2, left, top, height = height)

        left= Inches(0)
        top = Inches(1)
        height = Inches(1) 
        width = Inches(6)
        findingTxtBox = slideFinding.shapes.add_textbox(left,top ,width,height)

        findingTf = findingTxtBox.text_frame
        findingTf.text = ""

        findingParagraph = findingTf.add_paragraph()
        findingParagraph.text = ("Finding-"+ finding['hostName'])
        findingParagraph.font.size= Pt(30)


    #-----START HISTROGRAM SLIDE-----#

    slideHisto = ppt.slides.add_slide(ppt.slide_layouts[5])

         
    left = Inches(.5)
    top = Inches(0)    
    height = Inches(1)  
    
    #Add image
    pic = slideHisto.shapes.add_picture(img_path, left, top, height = height)

    left = Inches(8)
    top = Inches(0)    
    height = Inches(1.4)  
    
    pic = slideHisto.shapes.add_picture(img_path2, left, top, height = height)
    left= Inches(0)
    top = Inches(1)
    height = Inches(1) 
    width = Inches(6)
    findingTxtBox = slideHisto.shapes.add_textbox(left,top ,width,height)
    findingTf = findingTxtBox.text_frame
    findingTf.text = ""
    findingParagraph = findingTf.add_paragraph()
    findingParagraph.text = "Findings Histogram:"
    findingParagraph.font.size= Pt(30)

    info = 0
    veryLow = 0
    low = 0
    medium = 0
    high = 0
    veryHigh= 0


    for x in range(len(finding_json)):
        finding = finding_json[x]
        if finding["findingRisk"] == 'INFO':
            info += 1
        elif finding["findingRisk"] == 'VL':
            veryLow += 1
        elif finding["findingRisk"] == 'L':
            low += 1
        elif finding["findingRisk"] == 'M':
            medium += 1
        elif finding["findingRisk"] == 'H':
            high += 1
        else:
            veryHigh += 1

    
    chart_data = CategoryChartData()
    chart_data.categories = ['INFO', 'VERY LOW', 'LOW', 'MEDIUM', 'HIGH', 'VERY HIGH']
    chart_data.add_series('Series 1', (info, veryLow, low, medium, high, veryHigh))

    

    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(5)
    slideHisto.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

    ppt.save("../src/reports/ERB.pptx")
    return "OK"


@app.route("/generatefinalreport", methods=["POST"])
def generatefinalreport():

    document = Document()  # Report Object
    # Formatting
    TStyle = document.styles["Normal"]
    # Font
    TStyle.font.name = "Calibri (Body)"
    # Font Size
    TStyle.font.size = Pt(12)
    # How to add a Picture
    document.add_picture(
        "../src/assets/logo.png",
    )

    sentence = document.add_paragraph()

    sentence.add_run(
        "Combat Capabilities Development Command (CCDC) Data & Analysis Center (DAC) Enter System Name Enter Event Type (e.g., CVPA, CVI, VoF, etc) Report"
    ).bold = True

    # Specify all the analyst in the event.
    sentence.add_run("by ").bold = True
    # End of the page description
    sentence.add_run(
        "Classified by: Enter Lead Analyst Name\nDerived from: Enter Title of System's Security Classification Guide\nDeclassify on: Enter Declassification Date (e.g., 04/20/2040)"
    ).bold = True
    document.add_page_break()

    # Connect to the finding collection
    myclient = pymongo.MongoClient("mongodb://localhost:27017/")
    mydb = myclient["FRIC"]  # Database name
    mycollection = mydb["finding"]  # Collection Name

    finding_json = []

    # Start of Finding
    for e in mycollection.find():
        finding_json.append(
            {
                "id": e["id"],
                "hostName": e["Host_Name"],
                "ip_port": e["IP_Port"],
                "description": e["Description"],
                "longDescription": e["Long_Description"],
                "findingStatus": e["Finding_Status"],
                "findingType": e["Finding_Type"],
                "findingClassification": e["Finding_Classification"],
                "findingSystem": e["Finding_System"],
                "findingTask": e["Finding_Task"],
                "findingSubtask": e["Finding_Subtask"],
                "relatedFindings": e["Related_Findings"],
                "findingConfidentiality": e["Finding_Confidentiality"],
                "findingIntegrity": e["Finding_Integrity"],
                "findingAvailability": e["Finding_Availability"],
                "findingAnalyst": e["Finding_Analyst"],
                "findingCollaborators": e["Finding_Collaborators"],
                "findingPosture": e["Finding_Posture"],
                "mitigationDesc": e["Mitigation_Desc"],
                "mitigationLongDesc": e["Mitigation_Long_Desc"],
                "threatRelevence": e["Threat_Relevence"],
                "countermeasure": e["Countermeasure"],
                "impactDesc": e["Impact_Desc"],
                "impactLevel": e["Impact_Level"],
                "severityCategoryScore": e["Severity_Score"],
                "vulnerabilityScore": e["Vulnerability_Score"],
                "quantitativeScore": e["Quantitative_Score"],
                "findingRisk": e["Finding_Risk"],
                "findingLikelihood": e["Finding_Likelihood"],
                "findingCFIS": e["Finding_CFIS"],
                "findingIFIS": e["Finding_IFIS"],
                "findingAFIS": e["Finding_AFIS"],
                "impactScore": e["Impact_Score"],
                "findingFiles": e["Finding_Files"],
                "severityCategoryCode": e["Severity_Category_Code"],
                "systemID": e["System_ID"],
                "taskID": e["Task_ID"],
                "subtaskID": e["Subtask_ID"],
            }
        )

    # Create the tables
    for x in range(len(finding_json)):
        tableDescription = document.add_paragraph()
        finding = finding_json[x]
        tableDescription.add_run(
            "Table " + str(x + 1) + " describes the " + finding["description"] + " "
        )
        tableDescription.add_run(
            "\nTable " + str(x + 1) + ". " + finding["description"]
        ).bold = True
        table = document.add_table(
            rows=1, cols=8
        )  # Specify Rows and Columns for the Table
        table.autofit = True
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = "ID"
        hdr_cells[1].text = finding["id"]
        hdr_cells[3].text = "Impact Score"
        hdr_cells[4].text = str(finding["impactScore"])
        hdr_cells[5].text = "Status"
        hdr_cells[6].text = finding["findingStatus"]
        hdr_cells[7].text = "Posture"
        row_cells = table.add_row().cells  # Add Row
        row_cells[0].text = "Host Names"
        row_cells[2].text = "IP:Port"
        row_cells[3].text = "CAT"
        row_cells[4].text = finding["severityCategoryCode"]
        row_cells[5].text = "Likelihood"
        row_cells[6].text = finding["findingLikelihood"]
        row_cells[7].text = finding["findingPosture"]
        row_cells = table.add_row().cells  # Add Row
        row_cells[3].text = "CAT Score"
        row_cells[4].text = str(finding["severityCategoryScore"])
        row_cells[5].text = "Impact"
        row_cells[6].text = finding["impactLevel"]
        row_cells = table.add_row().cells  # Add Row
        row_cells[3].text = "VS-Score"
        row_cells[4].text = str(finding["vulnerabilityScore"])
        row_cells[5].text = "Risk"
        row_cells[6].text = finding["findingRisk"]
        row_cells = table.add_row().cells  # Add Row

    document.save("../src/reports/finalreport.docx")
    return "OK"
